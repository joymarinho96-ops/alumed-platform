import os
import fitz  # PyMuPDF
import re
import json
import hashlib
from pptx import Presentation

# ==========================================
# 1. CONFIGURACIÓN Y FUENTES
# ==========================================
CONFIG = {
    "rutaPrincipal": os.path.abspath("pdf"),
    "materiasPermitidas": [
        "anatomia-a",
        "anatomia-b",
        "anatomia-c",
        "biologia",
        "histologia-embriologia"
    ],
    "buscarRecursivamente": True,
    "aceptarMateriasNoListadas": False,
    "usarFuentesLegadas": False
}

MAPAS_TP = {
    "Biología": {
        "TP1": {"nombre": "Introducción a la Biología & Organización Celular", "keywords": ["procariota", "eucariota", "célula", "organela", "virus", "viroide", "prion"]},
        "TP2": {"nombre": "Componentes Químicos I (Agua y Pequeñas Moléculas)", "keywords": ["agua", "puente de hidrógeno", "ph", "buffer", "tampón", "ion", "hidrofílico", "hidrofóbico"]},
        "TP3": {"nombre": "Componentes Químicos II (Lípidos y Macromoléculas)", "keywords": ["lípido", "proteína", "carbohidrato", "aminoácido", "enlace peptídico", "glúcido", "ácido graso", "triglicérido", "fosfolípido"]},
        "TP4": {"nombre": "Bioenergética y Cinética Enzimática", "keywords": ["enzima", "catálisis", "sitio activo", "km", "vmax", "inhibidor", "alostérico", "atp", "energía libre", "termodinámica"]},
        "TP5": {"nombre": "Mecanismos Genéticos Básicos I", "keywords": ["adn", "arn", "nucleótido", "replicación", "polimerasa", "horquilla", "codón", "genes"]},
        "TP6": {"nombre": "Mecanismos Genéticos Básicos II (Expresión y Regulación)", "keywords": ["transcripción", "traducción", "promotor", "operón", "splicing", "intrón", "exón", "ribosoma", "arnt", "arnm"]},
        "TP7": {"nombre": "Membrana Plasmática y Transporte (Ósmosis/Osmolaridad)", "keywords": ["membrana", "transporte", "difusión", "ósmosis", "osmolaridad", "bomba", "na+/k+", "canales", "acuaporina", "bicapa"]},
        "TP8": {"nombre": "Metabolismo Intermedio y Mitocondrias", "keywords": ["mitocondria", "glucólisis", "krebs", "cadena respiratoria", "fosforilación oxidativa", "piruvato", "atp sintasa"]},
        "TP9": {"nombre": "Membranas Internas I (RE y Golgi)", "keywords": ["retículo", "endoplásmico", "golgi", "glicosilación", "secreción", "vesícula", "somatico", "rugoso", "liso"]},
        "TP10": {"nombre": "Membranas Internas II (Lisosomas y Peroxisomas)", "keywords": ["lisosoma", "peroxisoma", "endosoma", "fagocitosis", "autofagia", "hidrolasa", "catalasa", "digestión celular"]},
        "TP11": {"nombre": "El Núcleo Celular", "keywords": ["núcleo", "envoltura nuclear", "poro nuclear", "nucleólo", "cromatina", "heterocromatina", "eucromatina", "histona"]},
        "TP12": {"nombre": "Citoesqueleto, Adhesión y Matriz Extracelular", "keywords": ["citoesqueleto", "microtúbulo", "microfilamento", "filamento intermedio", "actina", "tubulina", "cilio", "flagelo", "matriz extracelular", "colágeno"]},
        "TP13": {"nombre": "Ciclo Celular y Apoptosis", "keywords": ["ciclo celular", "g1", "g2", "s", "p53", "cdk", "ciclina", "checkpoint", "punto de control", "apoptosis", "caspasa"]},
        "TP14": {"nombre": "Mecanismos de División Celular (Mitosis y Meiosis)", "keywords": ["mitosis", "meiosis", "profase", "metafase", "anafase", "telofase", "huso", "crossing over", "recombinación", "gameto"]},
        "TP15": {"nombre": "Transmisión del Material Genético", "keywords": ["mendel", "herencia", "alelo", "genotipo", "fenotipo", "dominante", "recesivo", "ligado al sexo"]},
        "TP16": {"nombre": "Las Células en su Contexto Social (Comunicación)", "keywords": ["comunicación celular", "receptor", "ligando", "segundo mensajero", "ampc", "quinasa", "transducción", "señalización"]}
    },
    "Histología y Embriología": {
        "TP1": {"nombre": "Técnica Histológica", "keywords": ["técnica", "fijación", "hematoxilina", "eosina", "tinción", "microscopio", "corte", "parafina"]},
        "TP2": {"nombre": "Tejido Epitelial", "keywords": ["epitelio", "revestimiento", "glandular", "cilio", "microvellosidad", "lámina basal", "unión ocluyente", "desmosoma"]},
        "TP3": {"nombre": "Tejido Conectivo y Adiposo", "keywords": ["conectivo", "conjuntivo", "fibroblasto", "colágeno", "matriz", "adipocito", "grasa", "laxo", "denso"]},
        "TP4": {"nombre": "Tejido Cartilaginoso y Óseo", "keywords": ["hueso", "cartílago", "osteocito", "osteoblasto", "osteoclasto", "condrocito", "osteona", "havers", "trabécula"]},
        "TP5": {"nombre": "Sangre y Hematopoyesis", "keywords": ["sangre", "glóbulo", "eritrocito", "leucocito", "plaqueta", "médula ósea", "hematopoyesis", "neutrófilo", "linfocito"]},
        "TP6": {"nombre": "Tejido Nervioso y Sistema Nervioso", "keywords": ["neurona", "axón", "dendrita", "glía", "astroctio", "oligodendrocito", "mielina", "sinapsis", "sustancia gris", "sustancia blanca"]},
        "TP7": {"nombre": "Tejido Muscular", "keywords": ["músculo", "estriado", "liso", "cardíaco", "sarcómero", "miocito", "miofilamento", "disco intercalar"]},
        "TP8": {"nombre": "Sistema Cardiovascular", "keywords": ["vaso", "arteria", "vena", "capilar", "endotelio", "corazón", "miocardio", "endocardio"]},
        "TP9": {"nombre": "Tejido Linfoide y Sistema Linfoide", "keywords": ["linfático", "ganglio", "bazo", "timo", "amígdala", "folículo linfoide"]},
        "TP10": {"nombre": "Embriología: Fecundación y Primeras Semanas", "keywords": ["fecundación", "zocota", "mórula", "blastocisto", "implantación", "trofoblasto", "embrioblasto", "segregación"]},
        "TP11": {"nombre": "Embriología: 3ª y 4ª Semana (Gastrulación/Neurulación)", "keywords": ["gastrulación", "ectodermo", "mesodermo", "endodermo", "notocorda", "tubo neural", "somita", "neurulación"]},
        "TP12": {"nombre": "Embriología: Placenta y Anexos Embrionarios", "keywords": ["placenta", "corion", "amnios", "cordón umbilical", "vellosidad coriónica", "saco vitelino"]}
    },
    "Anatomía": {
        "TP1": {"nombre": "Huesos del Miembro Superior", "keywords": ["escápula", "clavícula", "húmero", "cúbito", "radio", "carpo", "metacarpo", "falange"]},
        "TP2": {"nombre": "Huesos del Miembro Inferior", "keywords": ["coxal", "fémur", "rótula", "tibia", "peroné", "tarso", "metatarso"]},
        "TP3": {"nombre": "Columna Vertebral y Tórax", "keywords": ["vértebra", "cervical", "torácica", "lumbar", "sacro", "costilla", "esternón"]},
        "TP4": {"nombre": "Cráneo (Bóveda y Base)", "keywords": ["frontal", "parietal", "occipital", "temporal", "esfenoides", "etmoides", "agujero", "fosa cranial"]},
        "TP5": {"nombre": "Macizo Facial", "keywords": ["maxilar", "cigomático", "nasal", "vómer", "mandíbula", "orbita"]},
        "TP6": {"nombre": "Artrología", "keywords": ["articulación", "sinovial", "ligamento", "cápsula", "menisco", "diartrosis"]},
        "TP7": {"nombre": "Miología Miembro Superior", "keywords": ["biceps", "triceps", "deltoides", "pectoral", "manguito rotador", "pronador", "supinador"]},
        "TP8": {"nombre": "Miología Miembro Inferior", "keywords": ["cuádriceps", "isquiotibial", "glúteo", "gastrocnemio", "sóleo", "aductor"]},
        "TP9": {"nombre": "Miología del Tronco", "keywords": ["recto abdominal", "oblicuo", "diafragma", "dorsal ancho", "trapecio"]},
        "TP10": {"nombre": "Miología Cabeza y Cuello", "keywords": ["estrenocleidomastoideo", "masetero", "temporal", "mimica", "suprahioideo"]}
    }
}

# ==========================================
# 2. RESOLUCIÓN DE METADATOS DE RUTA
# ==========================================
def resolver_metadatos_ruta(rel_path):
    parts = rel_path.replace("\\", "/").lower().split("/")
    materia_str = "Biología"
    catedra_str = "Única"
    modalidad_str = "CHOICE"
    parcial_str = "1"

    for p in parts:
        if p == "anatomia-a":
            materia_str, catedra_str = "Anatomía", "A"
        elif p == "anatomia-b":
            materia_str, catedra_str = "Anatomía", "B"
        elif p == "anatomia-c":
            materia_str, catedra_str = "Anatomía", "C"
        elif p == "biologia":
            materia_str, catedra_str = "Biología", "Única"
        elif p == "histologia-embriologia":
            materia_str, catedra_str = "Histología y Embriología", "Única"
            
        if p == "oral":
            modalidad_str = "ORAL"
        elif p == "pinches":
            modalidad_str = "PINCHE_STATION"
        elif p == "multiple-choice":
            modalidad_str = "CHOICE"
            
        if "parcial-1" in p:
            parcial_str = "1"
        elif "parcial-2" in p:
            parcial_str = "2"

    return materia_str, catedra_str, modalidad_str, parcial_str

def limpiar_texto(texto):
    if not texto:
        return ""
    patterns = [
        r"(?i)universidad\s+nacional\s+de\s+la\s+plata",
        r"(?i)facultad\s+de\s+ciencias\s+médicas",
        r"(?i)f\.?c\.?m\.?\s*-?\s*u\.?n\.?l\.?p\.?",
        r"(?i)cátedra\s+[a-c]?",
        r"(?i)autor\s+responsable:?.*",
        r"(?i)diseñador:?.*",
        r"(?i)editor:?.*",
        r"(?i)página\s+\d+(\s+de\s+\d+)?",
        r"(?i)figura\s+\d+:?.*",
        r"(?i)alumed\s+instituto.*",
        r"(?i)profe\s+joyce.*",
        r"^\s*\d+[\.\-\)]\s*"
    ]
    cl = texto
    for p in patterns:
        cl = re.sub(p, "", cl)
    return re.sub(r"\s+", " ", cl).strip()

def clasificar_pedagogicamente(item, materia):
    texto = ((item.get("pregunta", "") + " " + item.get("explicacion", "") + " " + json.dumps(item.get("opciones", [])))).lower()
    mapa = MAPAS_TP.get(materia, MAPAS_TP["Biología"])
    
    best_tp = "TP1"
    max_score = 0
    kw_match = []
    
    for tp_id, tp_info in mapa.items():
        score = 0
        matches = []
        for kw in tp_info["keywords"]:
            if kw in texto:
                score += 1
                matches.append(kw)
        if score > max_score:
            max_score = score
            best_tp = tp_id
            kw_match = matches
            
    confianza = "alta" if max_score >= 3 else ("media" if max_score >= 1 else "baja")
    
    item["tpPrincipal"] = best_tp
    item["tpId"] = best_tp
    item["tema"] = mapa.get(best_tp, {}).get("nombre", "General")
    item["subtema"] = kw_match[0].capitalize() if kw_match else "Conceptos Básicos"
    item["conceptosClave"] = kw_match if kw_match else ["general"]
    item["tpRelacionados"] = [t for t in mapa.keys() if t != best_tp][:2]
    item["justificacionClasificacion"] = f"Anclado a {best_tp} por términos temáticos clave: {', '.join(kw_match[:3]) if kw_match else 'revisión conceptual'}."
    item["confianzaClasificacion"] = confianza
    item["estadoClasificacion"] = "clasificado" if confianza in ["alta", "media"] else "pendiente_revision"
    return item

# ==========================================
# 3. EXTRACTOR POR TIPO DE ARCHIVO
# ==========================================
def extraer_de_pdf(filepath, rel_path):
    materia, catedra, modalidad, parcial = resolver_metadatos_ruta(rel_path)
    filename = os.path.basename(filepath)
    preguntas = []
    
    doc = fitz.open(filepath)
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text")
        images = page.get_images(full=True)
        
        blocks = re.split(r'\n(?=\d+[\.\-\)]\s*)', text)
        for b in blocks:
            lines = [l.strip() for l in b.strip().split('\n') if l.strip()]
            if not lines or len(lines[0]) < 10:
                continue
                
            enunciado = limpiar_texto(lines[0])
            if not enunciado:
                continue
                
            opciones = []
            correcta = 0
            for l in lines[1:]:
                m_opt = re.match(r'^\s*([a-dA-D])[\.\-\)]\s*(.*)', l)
                if m_opt:
                    opt_t = m_opt.group(2).strip()
                    is_corr = '*' in l or 'CORRECTA' in l.upper() or '[OK]' in l
                    opciones.append({
                        "texto": re.sub(r'[\*[OK]]', '', opt_t).strip(),
                        "explicacion": "Opción validada de examen oficial."
                    })
                    if is_corr:
                        correcta = len(opciones) - 1
                        
            if len(opciones) == 0 and modalidad != "CHOICE":
                opciones = [{"texto": "Respuesta oral basada en bibliografía oficial.", "explicacion": ""}]
                
            if len(opciones) < 2 and modalidad == "CHOICE":
                continue
                
            q_id = "EXC-" + hashlib.md5((filename + str(page_num) + enunciado).encode('utf-8')).hexdigest()[:10]
            
            q_item = {
                "id": q_id,
                "materia": materia,
                "catedra": catedra,
                "parcial": parcial,
                "modalidad": modalidad,
                "pregunta": enunciado,
                "opciones": opciones,
                "correcta": correcta,
                "explicacion": "Respuesta basada en materiales ALUMED.",
                "fuenteInterna": rel_path,
                "textoContextualInterno": lines[0]
            }
            
            clasificar_pedagogicamente(q_item, materia)
            preguntas.append(q_item)
            
    doc.close()
    return preguntas

def extraer_de_pptx(filepath, rel_path):
    materia, catedra, modalidad, parcial = resolver_metadatos_ruta(rel_path)
    filename = os.path.basename(filepath)
    preguntas = []
    
    prs = Presentation(filepath)
    for idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
                
        full_text = "\n".join(slide_text).strip()
        if not full_text:
            continue
            
        enunciado = limpiar_texto(full_text.split('\n')[0])
        if not enunciado:
            continue
            
        q_id = "PPT-" + hashlib.md5((filename + str(idx) + enunciado).encode('utf-8')).hexdigest()[:10]
        
        q_item = {
            "id": q_id,
            "materia": materia,
            "catedra": catedra,
            "parcial": parcial,
            "modalidad": modalidad,
            "pregunta": enunciado,
            "opciones": [{"texto": "Identificación de preparado / estación de Pinche.", "explicacion": ""}],
            "correcta": 0,
            "explicacion": "Estación de pinche / diapositiva evaluativa.",
            "fuenteInterna": rel_path
        }
        
        clasificar_pedagogicamente(q_item, materia)
        preguntas.append(q_item)
        
    return preguntas

# ==========================================
# 4. PROCESAMIENTO RECURSIVO Y MANIFIESTO
# ==========================================
def ejecutar_excavacion_definitiva():
    manifiesto_rutas = []
    todas_preguntas = []
    
    root_pdf = os.path.join(os.getcwd(), "pdf")
    print(f"=== INICIANDO EXCAVACIÓN DEFINITIVA EN {root_pdf} ===")
    
    for mat_dir in CONFIG["materiasPermitidas"]:
        mat_path = os.path.join(root_pdf, mat_dir)
        if not os.path.exists(mat_path):
            print(f"Ruta no encontrada: {mat_path}")
            continue
            
        for root, dirs, files in os.walk(mat_path):
            rel_dir = os.path.relpath(root, root_pdf)
            materia, catedra, modalidad, parcial = resolver_metadatos_ruta(rel_dir)
            
            pdf_count = 0
            ppt_count = 0
            img_count = 0
            qs_dir = []
            
            for file in files:
                full_f = os.path.join(root, file)
                rel_f = os.path.relpath(full_f, root_pdf)
                ext = os.path.splitext(file)[1].lower()
                
                if ext == ".pdf":
                    pdf_count += 1
                    qs = extraer_de_pdf(full_f, rel_f)
                    qs_dir.extend(qs)
                elif ext in [".ppt", ".pptx"]:
                    ppt_count += 1
                    qs = extraer_de_pptx(full_f, rel_f)
                    qs_dir.extend(qs)
                elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
                    img_count += 1
                    
            todas_preguntas.extend(qs_dir)
            
            manifiesto_rutas.append({
                "ruta": rel_dir,
                "materia": materia,
                "catedra": catedra,
                "modalidad": modalidad,
                "parcial": parcial,
                "archivosEncontrados": len(files),
                "pdfEncontrados": pdf_count,
                "presentacionesEncontradas": ppt_count,
                "imagenesEncontradas": img_count,
                "preguntasImportadas": len(qs_dir),
                "estado": "completo"
            })

    # Deduplicar e integrar en data.js
    with open("data.js", "r", encoding="utf-8") as f:
        content = f.read()
        
    first_b = content.find('{')
    last_b = content.rfind('}')
    banco = json.loads(content[first_b:last_b+1])
    
    existing = banco.get("choices", [])
    seen_texts = set(limpiar_texto(q.get("pregunta", "")) for q in existing)
    
    agregadas = 0
    for q in todas_preguntas:
        if q["pregunta"] not in seen_texts:
            existing.append(q)
            seen_texts.add(q["pregunta"])
            agregadas += 1
            
    banco["choices"] = existing
    
    out_code = "const bancoDados = " + json.dumps(banco, ensure_ascii=False, indent=2) + ";"
    with open("data.js", "w", encoding="utf-8") as f:
        f.write(out_code)
        
    # Guardar manifiesto
    with open("manifiesto_excavacion.json", "w", encoding="utf-8") as f:
        json.dump(manifiesto_rutas, f, ensure_ascii=False, indent=2)
        
    print(f"\n[OK] EXCAVACIÓN COMPLETADA EXITOSAMENTE.")
    print(f"Nuevas preguntas integradas: {agregadas}")
    print(f"Total general en data.js: {len(existing)}")

if __name__ == "__main__":
    ejecutar_excavacion_definitiva()
