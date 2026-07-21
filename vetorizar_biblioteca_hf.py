import os
import sys
import django
from pathlib import Path
import PyPDF2
from pptx import Presentation
from fastembed import TextEmbedding

# 1. Configurar o ambiente Django para podermos acessar o banco de dados (ProfeJoyChunk)
sys.path.append(str(Path(__file__).resolve().parent))
os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'alumed.settings')
django.setup()

from accounts.models import ProfeJoyChunk

def extract_text_from_pdf(pdf_path):
    """Lê o PDF e retorna o texto completo."""
    try:
        reader = PyPDF2.PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        print(f"Erro ao ler PDF {pdf_path}: {e}")
        return ""

def extract_text_from_pptx(pptx_path):
    """Lê a apresentação PowerPoint e retorna o texto completo."""
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Erro ao ler PPTX {pptx_path}: {e}")
        return ""

def chunk_text(text, chunk_size=2000, overlap=200):
    """Divide o texto em pedaços de 'chunk_size' caracteres com 'overlap' caracteres de sobreposição."""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - overlap
    return chunks

def run():
    biblioteca_dir = os.path.join(os.getcwd(), 'biblioteca_alumed')
    if not os.path.exists(biblioteca_dir):
        print(f"⚠️ Pasta não encontrada: {biblioteca_dir}")
        print("Crie a pasta 'biblioteca_alumed' e coloque seus PDFs dentro dela antes de rodar este script.")
        return

    # Limpar embeddings antigos para evitar duplicatas e também porque
    # o tamanho do vetor do OpenAI (1536) não bate com o do FastEmbed (384)
    print("🗑️ Limpando banco de dados antigo de PDFs...")
    ProfeJoyChunk.objects.filter(source_type='pdf').delete()

    print("🚀 Baixando/Carregando modelo de IA Open Source (HuggingFace/FastEmbed)...")
    print("Isso pode demorar alguns minutos na primeira vez para baixar os pesos da rede neural (~50MB).")
    embedding_model = TextEmbedding("BAAI/bge-small-en-v1.5")
    
    print("\n📚 Lendo PDFs da pasta biblioteca_alumed...\n")
    for root, dirs, files in os.walk(biblioteca_dir):
        for filename in files:
            if filename.lower().endswith(".pdf"):
                pdf_path = os.path.join(root, filename)
                print(f"📖 Processando: {filename}")
                
                text = extract_text_from_pdf(pdf_path)
                if not text:
                    print(f"   ⚠️ Nenhum texto extraído de {filename}. Pode ser um PDF de imagens.")
                    continue
                
                chunks = chunk_text(text)
                print(f"   🧩 Dividido em {len(chunks)} partes. Vetorizando...")
                
                # Gera as embeddings usando o modelo open source local (grátis!)
                embeddings_generator = embedding_model.embed(chunks)
                embeddings_list = list(embeddings_generator)
                
                # Salvar no banco de dados do Django
                for i, (chunk_text_data, emb) in enumerate(zip(chunks, embeddings_list)):
                    # O FastEmbed retorna numpy arrays, precisamos converter para lista de floats
                    emb_list = emb.tolist()
                    
                    ProfeJoyChunk.objects.create(
                        title=filename,
                        content=chunk_text_data,
                        embedding=emb_list,
                        source_type='pdf',
                        chunk_index=i,
                        subject="Geral",  # Você pode mudar de acordo com as pastas
                        year="1"
                    )
                
                print(f"   ✅ {filename} salvo com sucesso no 'cérebro' da Profe Joy!\n")
                
    print("🎉 VETORIZAÇÃO CONCLUÍDA! A Profe Joy agora sabe tudo sobre esses PDFs.")
    print("O motor Open Source está pronto para uso e sem cobrar um centavo de saldo. 💻⚡")

if __name__ == '__main__':
    run()
